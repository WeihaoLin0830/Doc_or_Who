from __future__ import annotations

import zipfile
from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation

from backend.types import ExtractionResult

UNSAFE_EXTENSIONS = {".docm", ".pptm", ".xlsm"}
MAX_ZIP_ENTRIES = 2000
MAX_UNCOMPRESSED_BYTES = 200 * 1024 * 1024


def validate_office_zip(path: Path) -> None:
    ext = path.suffix.lower()
    if ext in UNSAFE_EXTENSIONS:
        raise ValueError(f"Unsupported macro-enabled Office file: {ext}")
    if not zipfile.is_zipfile(path):
        raise ValueError("Invalid Office container")
    total_uncompressed = 0
    with zipfile.ZipFile(path) as archive:
        infos = archive.infolist()
        if len(infos) > MAX_ZIP_ENTRIES:
            raise ValueError("Office archive has too many entries")
        for info in infos:
            name = info.filename
            if name.startswith("/") or ".." in Path(name).parts:
                raise ValueError("Unsafe path inside Office archive")
            total_uncompressed += info.file_size
            if total_uncompressed > MAX_UNCOMPRESSED_BYTES:
                raise ValueError("Office archive is too large when uncompressed")


def extract_docx(path: Path) -> ExtractionResult:
    validate_office_zip(path)
    document = DocxDocument(str(path))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    props = document.core_properties
    metadata = {
        "title": props.title or None,
        "author": props.author or None,
        "subject": props.subject or None,
        "paragraph_count": len(paragraphs),
    }
    return ExtractionResult(path=path, text="\n\n".join(paragraphs), metadata=metadata)


def extract_pptx(path: Path) -> ExtractionResult:
    validate_office_zip(path)
    deck = Presentation(str(path))
    slides_text: list[str] = []
    for slide_number, slide in enumerate(deck.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and str(shape.text).strip():
                texts.append(str(shape.text).strip())
        if texts:
            slides_text.append(f"Slide {slide_number}\n" + "\n".join(texts))
    props = deck.core_properties
    metadata = {
        "title": props.title or None,
        "author": props.author or None,
        "subject": props.subject or None,
        "slide_count": len(deck.slides),
    }
    return ExtractionResult(path=path, text="\n\n".join(slides_text), metadata=metadata)
