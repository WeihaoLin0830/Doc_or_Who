from __future__ import annotations

from pathlib import Path

import fitz
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation


def write_text(path: Path, content: str) -> Path:
    path.write_text(content, encoding="utf-8")
    return path


def write_pdf(path: Path, pages: list[str]) -> Path:
    document = fitz.open()
    for text in pages:
        page = document.new_page()
        if text:
            page.insert_text((72, 72), text)
    document.save(path)
    document.close()
    return path


def write_blank_pdf(path: Path) -> Path:
    document = fitz.open()
    document.new_page()
    document.save(path)
    document.close()
    return path


def write_docx(path: Path, paragraphs: list[str]) -> Path:
    document = DocxDocument()
    for paragraph in paragraphs:
        document.add_paragraph(paragraph)
    document.core_properties.title = paragraphs[0]
    document.save(path)
    return path


def write_pptx(path: Path, slides: list[list[str]]) -> Path:
    presentation = Presentation()
    for slide_lines in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_lines[0]
        slide.placeholders[1].text = "\n".join(slide_lines[1:])
    presentation.core_properties.title = slides[0][0]
    presentation.save(path)
    return path


def write_xlsx(path: Path, rows: list[dict[str, object]], sheet_name: str = "Sheet1") -> Path:
    frame = pd.DataFrame(rows)
    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        frame.to_excel(writer, index=False, sheet_name=sheet_name)
    return path


def write_csv(path: Path, rows: list[dict[str, object]]) -> Path:
    frame = pd.DataFrame(rows)
    frame.to_csv(path, index=False)
    return path
